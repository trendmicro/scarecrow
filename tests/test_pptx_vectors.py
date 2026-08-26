"""Layer-1 tests for the PPTX vector library (Slice 5).

Deterministic, offline. Mirrors test_vector_library.py for PPTX:
  - all 10 PPTX vectors are registered, in order
  - each vector applied alone: its payload reaches the file, and clean removes it
    (visible content preserved; sidecar restored for text-value vectors)
  - --all applies the default (non-research-only) set and clean reverses it
  - the two encoding vectors are excluded from --all but stay in the library
  - inspect reports vector + location; sidecar written only for text-value vectors
"""
from __future__ import annotations

import zipfile

import pytest
from pptx import Presentation
from pptx.util import Inches

from scarecrow.formats import pptx as _pptx  # noqa: F401  (registers vectors)
from scarecrow import vector
from scarecrow.core.marker import Marker, write_sidecar, sidecar_path_for
from scarecrow.core.modes import Mode
from scarecrow.core.payloads import select_payload
from scarecrow.formats.pptx.reverse import inspect_pptx, clean_pptx

ALL = ["speaker_notes", "white_text", "tiny_font", "off_slide", "zero_size",
       "alt_text", "hidden_slide", "core_props", "unicode_tags", "zero_width"]
ENCODED = {"unicode_tags", "zero_width"}
TEXT_VALUE = {"speaker_notes", "core_props"}
PAYLOAD_MARK = "protected against automated ai"
VISIBLE = ["Q3 Report", "Revenue was 4.2M"]


def _make_original(tmp_path):
    p = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # "Title Only"
    slide.shapes.title.text = VISIBLE[0]
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    box.text_frame.text = VISIBLE[1]
    prs.save(str(p))
    return p


def _protect(orig, tmp_path, vectors: list[str]):
    out = tmp_path / "deck.scarecrow.pptx"
    prs = Presentation(str(orig))
    marker = Marker()
    payload = select_payload(Mode.REFUSE, None)
    for name in vectors:
        vector.get(name, "pptx").apply(prs, payload, marker)
    prs.save(str(out))
    write_sidecar(str(out), marker.sidecar_entries)
    return out


def _all_xml(path) -> str:
    z = zipfile.ZipFile(str(path))
    return "".join(z.read(n).decode("utf-8", "replace")
                   for n in z.namelist() if n.endswith(".xml"))


def _visible(path):
    prs = Presentation(str(path))
    texts = []
    for slide in prs.slides:
        if slide.element.get("show") == "0":
            continue  # a hidden slide is not "visible" content
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                texts.append(sh.text_frame.text.strip())
    return texts


def _decode_present(xml: str) -> bool:
    from_tags = "".join(
        chr(ord(c) - 0xE0000) if 0xE0000 <= ord(c) <= 0xE007F else "" for c in xml)
    if "protected" in from_tags.lower():
        return True
    return (xml.count("​") + xml.count("‌")) > 100


def test_all_pptx_vectors_registered():
    names = [v.name for v in vector.all_vectors("pptx")]
    assert names == ALL  # order matters for --all / list-vectors


@pytest.mark.parametrize("name", ALL)
def test_vector_payload_reaches_the_file(name, tmp_path):
    orig = _make_original(tmp_path)
    prot = _protect(orig, tmp_path, [name])
    xml = _all_xml(prot)
    if name in ENCODED:
        assert _decode_present(xml)
    else:
        assert (PAYLOAD_MARK in xml.lower()) or ("owner-protected" in xml.lower())


@pytest.mark.parametrize("name", ALL)
def test_vector_clean_restores_original(name, tmp_path):
    orig = _make_original(tmp_path)
    prot = _protect(orig, tmp_path, [name])
    cleaned = tmp_path / "cleaned.pptx"
    clean_pptx(str(prot), str(cleaned))
    # visible content preserved
    assert set(VISIBLE).issubset(set(_visible(cleaned)))
    # no payload left anywhere
    xml = _all_xml(cleaned).lower()
    assert PAYLOAD_MARK not in xml
    assert "owner-protected" not in xml
    if name in ENCODED:
        assert not _decode_present(_all_xml(cleaned))


def test_all_applies_defaults_and_clean_reverses(tmp_path):
    orig = _make_original(tmp_path)
    defaults = [v.name for v in vector.default_vectors("pptx")]
    prot = _protect(orig, tmp_path, defaults)
    found = {i.vector for i in inspect_pptx(str(prot))}
    assert found == set(defaults)
    cleaned = tmp_path / "cleaned.pptx"
    removed = clean_pptx(str(prot), str(cleaned))
    assert removed == len(defaults)
    assert inspect_pptx(str(cleaned)) == []
    assert set(VISIBLE).issubset(set(_visible(cleaned)))


def test_default_set_excludes_encoding_vectors():
    default_names = [v.name for v in vector.default_vectors("pptx")]
    assert set(ENCODED).isdisjoint(default_names)          # excluded
    assert set(default_names) == set(ALL) - ENCODED        # everything else kept
    assert set(ENCODED).issubset(v.name for v in vector.all_vectors("pptx"))  # still in library


def test_sidecar_written_only_for_text_value_vectors(tmp_path):
    orig = _make_original(tmp_path)
    # structural-only: no sidecar
    prot1 = _protect(orig, tmp_path, ["off_slide"])
    assert not sidecar_path_for(str(prot1)).exists()
    # includes a text-value vector: sidecar exists
    prot2 = _protect(orig, tmp_path, ["off_slide", "speaker_notes"])
    assert sidecar_path_for(str(prot2)).exists()


def test_inspect_reports_vector_and_location(tmp_path):
    orig = _make_original(tmp_path)
    prot = _protect(orig, tmp_path, ["speaker_notes", "off_slide"])
    found = {i.vector: i.where for i in inspect_pptx(str(prot))}
    assert found["off_slide"] == "slide:0"
    assert found["speaker_notes"].startswith("sidecar:")
